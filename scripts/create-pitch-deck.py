#!/usr/bin/env python3.13
"""Generate the HackCanton S2 pitch deck for Cognivern."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Palette: Midnight Executive
NAVY = RGBColor(11, 15, 47)
ICE_BLUE = RGBColor(202, 220, 252)
WHITE = RGBColor(255, 255, 255)
CORAL = RGBColor(249, 97, 103)
MINT = RGBColor(2, 195, 154)
SLATE = RGBColor(54, 69, 79)

OUT_PATH = "docs/pitch-deck.pptx"


def add_background(slide, color=NAVY):
    """Fill the slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return box


def add_shape_text(slide, left, top, width, height, text, fill_color=None,
                   font_size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
                   paragraphs=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.color.rgb = ICE_BLUE
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)

    parts = paragraphs if paragraphs else [(text, font_size, bold)]
    tf.clear()
    for i, part in enumerate(parts):
        if len(part) == 3:
            txt, size, is_bold = part
            col = color
        else:
            txt, size, is_bold, col = part
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.bold = is_bold
        p.font.color.rgb = col
        p.font.name = "Calibri"
        p.alignment = align
    return shape


def title_slide(prs):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide)

    # Accent bar at top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CORAL
    bar.line.fill.background()

    add_textbox(slide, Inches(0.75), Inches(1.9), Inches(11.8), Inches(1.4),
                "Cognivern", font_size=66, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.75), Inches(3.2), Inches(11.8), Inches(0.9),
                "Private sealed-bid RFP / OTC vendor selection on Canton Network",
                font_size=26, color=ICE_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.75), Inches(4.25), Inches(11.8), Inches(0.7),
                "Track 1 · Private DeFi & Capital Markets  |  Team thisyearnofear",
                font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.75), Inches(5.5), Inches(11.8), Inches(0.6),
                "Live demo: cognivern.vercel.app/sealed-bid",
                font_size=16, color=MINT, align=PP_ALIGN.CENTER)
    return slide


def problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.8),
                "The pain: information leakage kills institutional RFPs and OTC", font_size=36, bold=True, color=WHITE)

    cards = [
        ("Email & portals leak", "A leaked quote in an RFP can cost 5–15% in final price. Respondents who see a competitor’s price band the whole auction."),
        ("SaaS counterparty risk", "Unblinding is centralized in one provider — a new trust failure inside an outsourced workflow."),
        ("OTC desks freeze up", "Market makers won’t quote tight spreads if counterparties fear print-level leakage."),
    ]
    x = Inches(0.6)
    y = Inches(1.7)
    w = Inches(3.9)
    h = Inches(2.8)
    gap = Inches(0.35)
    for title, body in cards:
        add_shape_text(slide, x, y, w, h, "", fill_color=SLATE,
                       color=WHITE,
                       paragraphs=[(title, 20, True), (body, 14, False)])
        x += w + gap
    add_textbox(slide, Inches(0.6), Inches(4.75), Inches(12), Inches(0.6),
                "What institutions need: a sealed-bid primitive where bids exist, but only the counterparty can read them — and revealing the winner does not reveal every loser.",
                font_size=18, color=ICE_BLUE, align=PP_ALIGN.CENTER)
    return slide


def canton_answer_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.8),
                "Canton makes privacy structural, not cryptographic", font_size=36, bold=True, color=WHITE)

    # Flow diagram
    boxes = [
        ("SealedBidAuction", "Auctioneer + eligible bidders"),
        ("SubmitBid", "Each bidder creates a Bid contract"),
        ("CloseAndReveal", "One consuming choice"),
        ("AuctionResult", "Winner + winning amount public"),
    ]
    x = Inches(0.6)
    y = Inches(1.8)
    w = Inches(2.7)
    h = Inches(1.6)
    for title, body in boxes:
        add_shape_text(slide, x, y, w, h, "", fill_color=SLATE,
                       color=WHITE, align=PP_ALIGN.CENTER,
                       paragraphs=[(title, 18, True), (body, 13, False)])
        # arrow
        if x < Inches(9):
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + w + Inches(0.08), y + Inches(0.55), Inches(0.5), Inches(0.5))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = CORAL
            arrow.line.fill.background()
        x += w + Inches(0.66)

    add_textbox(slide, Inches(0.6), Inches(3.8), Inches(12), Inches(1.2),
                "The Bid template gives signatory rights to the bidder and observer rights only to the auctioneer. No encryption key ceremony. No threshold of losers. The reveal is one atomic ledger transaction that consumes every losing Bid in flight.",
                font_size=18, color=ICE_BLUE, align=PP_ALIGN.CENTER)
    return slide


def comparison_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.8),
                "Why Canton beats a public chain or FHE-only path", font_size=34, bold=True, color=WHITE)

    headers = ["Public chain / SaaS", "FHE only (Fhenix)", "Canton Network"]
    rows = [
        ["Bid amounts visible to all", "Amounts stay sealed", "Amounts stay sealed"],
        ["Reveal leaks every loser", "Needs threshold decryption by losers", "Winner revealed without decrypting losers"],
        ["Trust a SaaS operator", "Manager publishes after decrypt", "Trust Daml disclosure model + participant node"],
        ["Settlement is off-ledger", "Reveal requires multi-party ceremony", "Atomic reveal in one ledger transaction"],
    ]
    col_x = [Inches(0.6), Inches(4.5), Inches(8.4)]
    col_w = Inches(3.7)
    row_h = Inches(0.65)
    start_y = Inches(1.4)
    for i, h in enumerate(headers):
        add_shape_text(slide, col_x[i], start_y, col_w, row_h, h, fill_color=CORAL if i == 2 else SLATE,
                       font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for r, row in enumerate(rows):
        y = start_y + Inches(0.75) + r * (row_h + Inches(0.08))
        for i, cell in enumerate(row):
            fill = NAVY if i == 2 else None
            add_shape_text(slide, col_x[i], y, col_w, row_h, cell, fill_color=fill,
                           font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    return slide


def daml_model_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.5), Inches(0.45), Inches(12), Inches(0.7),
                "The model is 79 lines of Daml — no external dependencies", font_size=32, bold=True, color=WHITE)

    code = '''template Bid with
  bidder : Party
  manager : Party
  amountUsd : Decimal
  ...
where
  signatory bidder
  observer manager          -- ONLY the auctioneer may inspect

choice CloseAndReveal : ContractId AuctionResult
  controller manager
  do
    ... select winner ...
    mapA archive losingBids  -- consumed in flight
    create AuctionResult with ...'''
    code_card = add_shape_text(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(3.7), "",
                   fill_color=SLATE, color=WHITE, align=PP_ALIGN.LEFT,
                   paragraphs=[(line, 14, False) for line in code.splitlines()])
    for p in code_card.text_frame.paragraphs:
        p.font.name = "Consolas"

    add_textbox(slide, Inches(0.5), Inches(5.25), Inches(12.3), Inches(0.9),
                "Two lines do the privacy work: observer manager on Bid, and CloseAndReveal being a consuming choice. The rest is plumbing.",
                font_size=16, color=ICE_BLUE, align=PP_ALIGN.CENTER)
    return slide


def product_slide(prs, screenshot_path=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.7),
                "Live product: sealed-bid vendor selection with a party view", font_size=32, bold=True, color=WHITE)
    if screenshot_path:
        slide.shapes.add_picture(screenshot_path, Inches(0.6), Inches(1.35), width=Inches(12))
    else:
        add_shape_text(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(4.2),
                       "[UI screenshot will be inserted here]\n\ncognivern.vercel.app/sealed-bid\n\nVisitors toggle Auctioneer / Alice / Bob / Charlie and instantly see which bids each party can read on the Canton ledger.",
                       fill_color=SLATE, font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    return slide


def rigor_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.7),
                "Engineering rigor: privacy guarantees are machine-verifiable", font_size=34, bold=True, color=WHITE)

    stats = [
        ("31+", "Vitest integration & unit tests"),
        ("24", "TestSprite CLI tests on the live API"),
        ("Direct", "Ledger assertions per party role"),
        ("79", "Lines of Daml across 3 templates"),
    ]
    x = Inches(0.5)
    y = Inches(1.5)
    w = Inches(2.9)
    h = Inches(2.3)
    gap = Inches(0.35)
    for num, label in stats:
        add_shape_text(slide, x, y, w, h, "", fill_color=SLATE,
                       color=WHITE, align=PP_ALIGN.CENTER,
                       paragraphs=[(num, 48, True, MINT), (label, 14, False, WHITE)])
        x += w + gap

    add_textbox(slide, Inches(0.6), Inches(4.15), Inches(12), Inches(1.2),
                "Privacy invariants are asserted by querying the Daml JSON Ledger API directly as each party — not by trusting our backend cache. Legal and risk teams get an immutable, non-repudiable audit trail.",
                font_size=18, color=ICE_BLUE, align=PP_ALIGN.CENTER)
    return slide


def compliance_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.7),
                "Built for legal, risk, and audit sign-off", font_size=34, bold=True, color=WHITE)

    cards = [
        ("Role-based disclosure", "Daml signatory/observer enforces who sees what — auctioneer, bidder, and auditor roles are structurally separated."),
        ("Non-repudiation", "Every bid submission and reveal is recorded on-ledger with hash-signed events in the CRE run ledger."),
        ("Immutable audit trail", "Full ledger history + AuctionResult contract gives procurement and compliance a durable record."),
        ("No SaaS trust gap", "Privacy comes from the disclosure model on your participant node — not from trusting a central unblinding operator."),
    ]
    x = Inches(0.5)
    y = Inches(1.35)
    w = Inches(5.9)
    h = Inches(1.9)
    for i, (title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        cx = x + col * (w + Inches(0.4))
        cy = y + row * (h + Inches(0.25))
        add_shape_text(slide, cx, cy, w, h, "", fill_color=SLATE,
                       color=WHITE, align=PP_ALIGN.LEFT,
                       paragraphs=[(title, 22, True, ICE_BLUE), (body, 14, False, WHITE)])
    return slide


def deployment_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.7),
                "Deployment path: DevNet today, private participant tomorrow", font_size=32, bold=True, color=WHITE)

    steps = [
        ("Today", "Live on HackCanton S2 DevNet with on-ledger proof artifacts and automated lifecycle verification (pnpm canton:proof)."),
        ("Private participant", "Same backend is participant-agnostic — point CANTON_JSON_API_URL at your own Canton node or a hosted validator."),
        ("Mainnet-ready model", "79-line Daml package compiles with SDK 3.5.x; no external dependencies. SettlementLeg upgrade path for atomic value transfer."),
    ]
    y = Inches(1.45)
    for i, (title, body) in enumerate(steps):
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.1), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = CORAL
        circle.line.fill.background()
        add_textbox(slide, Inches(0.6), y + Inches(0.12), Inches(0.55), Inches(0.5),
                    str(i + 1), font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_shape_text(slide, Inches(1.35), y, Inches(11.4), Inches(1.35),
                       "", fill_color=SLATE, color=WHITE, align=PP_ALIGN.LEFT,
                       paragraphs=[(title, 18, True, ICE_BLUE), (body, 14, False, WHITE)])
        y += Inches(1.55)
    return slide


def track_fit_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.7),
                "We answer the HackCanton builder-test questions", font_size=32, bold=True, color=WHITE)

    items = [
        ("Parties", "Auctioneer + eligible-bidder list enforced by Daml signatory/observer."),
        ("Visibility", "Each Bid is visible to its bidder and the auctioneer only."),
        ("Approval", "Only the manager’s CloseAndReveal choice can archive bids and emit the result."),
        ("Atomic settlement", "Winner selection + archive of every losing Bid + AuctionResult in one tx."),
        ("Auditability", "Ledger history + hash-signed events in the CRE run ledger."),
        ("Why weaker on public chain", "Public chains leak amounts; FHE needs a ceremony losers won’t join."),
    ]
    x = Inches(0.5)
    y = Inches(1.35)
    w = Inches(5.9)
    h = Inches(1.15)
    for i, (title, body) in enumerate(items):
        col = i % 2
        row = i // 2
        cx = x + col * (w + Inches(0.4))
        cy = y + row * (h + Inches(0.15))
        add_shape_text(slide, cx, cy, w, h, "", fill_color=SLATE,
                       color=WHITE, align=PP_ALIGN.LEFT,
                       paragraphs=[(title, 16, True, CORAL), (body, 13, False, WHITE)])
    return slide


def use_cases_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.7),
                "Real-world applicability: one primitive, four institutional workflows", font_size=32, bold=True, color=WHITE)

    cards = [
        ("Private RFPs", "Legal retainers, security audits, capex procurement where pricing is competitively sensitive."),
        ("Confidential OTC", "Client RFP for a block; cleared price visible, unmade quotes never surface."),
        ("Credit allocation", "Lenders bid into an allocation round; losing yield targets stay private."),
        ("B2B blind auctions", "Drop-in replacement for procurement portals that centralize unblinding."),
    ]
    x = Inches(0.5)
    y = Inches(1.35)
    w = Inches(5.9)
    h = Inches(1.9)
    for i, (title, body) in enumerate(cards):
        col = i % 2
        row = i // 2
        cx = x + col * (w + Inches(0.4))
        cy = y + row * (h + Inches(0.25))
        add_shape_text(slide, cx, cy, w, h, "", fill_color=SLATE,
                       color=WHITE, align=PP_ALIGN.LEFT,
                       paragraphs=[(title, 22, True, ICE_BLUE), (body, 14, False, WHITE)])
    return slide


def roadmap_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.7),
                "Roadmap: from informational settlement to atomic value settlement", font_size=32, bold=True, color=WHITE)

    steps = [
        ("Today", "CloseAndReveal archives bids and emits AuctionResult with winning amount."),
        ("Next", "Add a SettlementLeg to the auction and exercise the escrowed asset’s Settle choice inside the same CloseAndReveal transaction."),
        ("Bounty lane", "CBTC (BitSafe) private OTC escrow or cETH (OnRails) private collateral — 50,000 CC bounty."),
    ]
    y = Inches(1.45)
    for i, (title, body) in enumerate(steps):
        # numbered circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.1), Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = CORAL
        circle.line.fill.background()
        add_textbox(slide, Inches(0.6), y + Inches(0.12), Inches(0.55), Inches(0.5),
                    str(i + 1), font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_shape_text(slide, Inches(1.35), y, Inches(11.4), Inches(1.35),
                       "", fill_color=SLATE, color=WHITE, align=PP_ALIGN.LEFT,
                       paragraphs=[(title, 18, True, ICE_BLUE), (body, 14, False, WHITE)])
        y += Inches(1.55)
    return slide


def team_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.7),
                "Why us: a production platform, not a weekend mock-up", font_size=34, bold=True, color=WHITE)

    bullets = [
        "Built on top of an already-deployed, tested governance API — real production logging, auth, and audit rails.",
        "Daml model compiles clean with SDK 3.5.x; backend is participant-agnostic.",
        "End-to-end Canton flow works today: create → submit → close → atomic reveal, with privacy verified by direct ledger queries.",
        "31+ automated tests plus direct ledger assertions keep the privacy guarantee machine-verifiable.",
    ]
    y = Inches(1.5)
    for b in bullets:
        # bullet circle
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.15), Inches(0.18), Inches(0.18))
        dot.fill.solid()
        dot.fill.fore_color.rgb = MINT
        dot.line.fill.background()
        add_textbox(slide, Inches(1.05), y, Inches(11.5), Inches(0.8), b, font_size=18, color=WHITE)
        y += Inches(0.95)
    return slide


def closing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_textbox(slide, Inches(0.6), Inches(1.8), Inches(12), Inches(1.2),
                "Prove privacy on-ledger", font_size=58, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(3.2), Inches(12), Inches(0.9),
                "Cognivern brings institutional RFP / OTC sealed-bid workflows to Canton Network — with sub-transaction privacy and atomic reveal that actually works.",
                font_size=22, color=ICE_BLUE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(4.5), Inches(12), Inches(0.6),
                "Live demo: cognivern.vercel.app/sealed-bid", font_size=24, bold=True, color=MINT, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(5.35), Inches(12), Inches(0.6),
                "github.com/thisyearnofear/cognivern", font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide(prs)
    problem_slide(prs)
    canton_answer_slide(prs)
    comparison_slide(prs)
    daml_model_slide(prs)
    product_slide(prs)
    rigor_slide(prs)
    compliance_slide(prs)
    track_fit_slide(prs)
    use_cases_slide(prs)
    deployment_slide(prs)
    roadmap_slide(prs)
    team_slide(prs)
    closing_slide(prs)

    prs.save(OUT_PATH)
    print(f"Pitch deck saved to {OUT_PATH}")


if __name__ == "__main__":
    main()
